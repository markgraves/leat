"""Reader for docx files using docx2python."""

from pathlib import Path
from typing import Union
from zipfile import BadZipFile

from pptx import Presentation

from . import BaseReader


class PPTXReader(BaseReader):
    """Reader for pptx file"""

    def __init__(self):
        """
        Create a reader for PPTX files

        Returns:

        """
        pass

    def read(self, stream):
        """
        Extract pptx text from stream

        Args:
          stream: Stream to read from

        Returns:
          Text read from the stream
        """
        try:
            return self.get_text(Presentation(stream))
        except BadZipFile:
            print(
                "WARNING:",
                "Cannot read PPTX due to unknown format issue",
            )
            return ""

    def read_file(self, filename: Union[str, Path]):
        """
        Read text from pptx file

        Args:
          filename: Path | str: Filename to read from

        Returns:
          Text read from the file
        """
        try:
            return self.get_text(Presentation(filename))
        except BadZipFile:
            print(
                "WARNING:",
                "Cannot read PPTX due to unknown format issue:",
                filename,
            )
            return ""

    def get_text(self, prs: Presentation):
        """
        Read text from pptx presentation

        Args:
          filename: prs | Presentation: Presentation to read from

        Returns:
          Text read from the presentation
        """

        result = []
        for slide in prs.slides:
            slide_result = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_result.append(shape.text.strip())
            if slide_result:
                result.append("\n".join(slide_result))
        return "\n".join(result)
